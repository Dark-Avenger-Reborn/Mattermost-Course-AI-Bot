"""
rag/ingestor.py — Ingest course material into ChromaDB.

Pipeline:
  file → extract text → chunk → embed (bge-m3) → upsert into ChromaDB

Supported formats: PDF, PPTX, DOCX, XLSX/XLSM, MD, TXT, and most other
plain-text formats (py, html, csv, json, etc.)

Run via:  python scripts/ingest.py
Or call:  asyncio.run(ingest_all())
"""

import asyncio
import hashlib
import logging
from html.parser import HTMLParser
from pathlib import Path

import chromadb
import fitz  # pymupdf
from pptx import Presentation
from tqdm import tqdm

import config
from llm.client import embed

logger = logging.getLogger(__name__)

# Max files being embedded concurrently — keep low to avoid 504s on the gateway
_EMBED_SEMAPHORE = asyncio.Semaphore(2)

# File types to skip entirely (binary, media, etc.)
_SKIP_SUFFIXES = {
    ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".svg", ".ico", ".webp",
    ".mp3", ".mp4", ".wav", ".avi", ".mov", ".mkv",
    ".zip", ".tar", ".gz", ".rar", ".7z",
    ".exe", ".dll", ".so", ".dylib",
    ".pyc", ".pyo", "__pycache__",
    ".DS_Store", ".gitignore", ".git",
    ".db", ".sqlite", ".sqlite3",
}


# ── ChromaDB setup ──────────────────────────────────────────────────────────

def get_collection() -> chromadb.Collection:
    client = chromadb.PersistentClient(path=config.CHROMA_PATH)
    return client.get_or_create_collection(
        name="course_material",
        metadata={"hnsw:space": "cosine"},
    )


def reset_collection() -> None:
    """Delete the course_material collection so the next ingest recreates it fresh."""
    client = chromadb.PersistentClient(path=config.CHROMA_PATH)

    try:
        client.delete_collection(name="course_material")
        logger.info("Deleted ChromaDB collection course_material")
    except Exception:
        # The collection may not exist yet; treat that as already cleared.
        logger.info("ChromaDB collection course_material was already absent")


async def ensure_collection_compatible(collection: chromadb.Collection) -> chromadb.Collection:
    """Recreate the collection if its stored vector dimension differs from the active embedder."""
    try:
        sample = collection.get(include=["embeddings"], limit=1)
    except Exception as e:
        logger.warning("Could not inspect existing ChromaDB collection: %s", e)
        return collection

    embeddings = sample.get("embeddings")
    if embeddings is None or len(embeddings) == 0:
        return collection

    existing_dim = len(embeddings[0])
    current_dim = len((await embed(["dimension probe"]))[0])

    if existing_dim == current_dim:
        return collection

    logger.warning(
        "ChromaDB collection dimension %d does not match current embeddings %d; recreating collection",
        existing_dim,
        current_dim,
    )
    reset_collection()
    return get_collection()


# ── Text extraction ─────────────────────────────────────────────────────────

def extract_pdf(path: Path) -> list[dict]:
    doc = fitz.open(str(path))
    pages = []
    for i, page in enumerate(doc):
        text = page.get_text("text").strip()
        if text:
            pages.append({"page": i + 1, "text": text})
    doc.close()
    return pages


def extract_pptx(path: Path) -> list[dict]:
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
        text = "\n".join(parts).strip()
        if text:
            slides.append({"page": i + 1, "text": text})
    return slides


def extract_docx(path: Path) -> list[dict]:
    try:
        from docx import Document
        doc = Document(str(path))
        parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        text = "\n".join(parts).strip()
        return [{"page": 1, "text": text}] if text else []
    except ImportError:
        logger.warning("python-docx not installed; skipping %s. Run: pip install python-docx", path.name)
        return []


def extract_xlsx(path: Path) -> list[dict]:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(path), read_only=True, data_only=True)
        pages = []
        for si, name in enumerate(wb.sheetnames, 1):
            ws = wb[name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                if cells:
                    rows.append("\t".join(cells))
            text = "\n".join(rows).strip()
            if text:
                pages.append({"page": si, "text": f"[Sheet: {name}]\n{text}"})
        wb.close()
        return pages
    except ImportError:
        logger.warning("openpyxl not installed; skipping %s. Run: pip install openpyxl", path.name)
        return []


class _VisibleTextParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self._skip_depth = 0
        self._parts: list[str] = []

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag in {"script", "style"}:
            self._skip_depth += 1

    def handle_endtag(self, tag: str) -> None:
        if tag in {"script", "style"} and self._skip_depth:
            self._skip_depth -= 1

    def handle_data(self, data: str) -> None:
        if self._skip_depth:
            return

        text = data.strip()
        if text:
            self._parts.append(text)

    def get_text(self) -> str:
        return " ".join(self._parts).strip()


def extract_html(path: Path) -> list[dict]:
    try:
        raw = path.read_text(encoding="utf-8", errors="ignore")
    except Exception as e:
        logger.warning("Could not read %s as HTML: %s", path.name, e)
        return []

    parser = _VisibleTextParser()
    parser.feed(raw)
    text = parser.get_text()
    return [{"page": 1, "text": text}] if text else []


def extract_text(path: Path) -> list[dict]:
    """Generic fallback: read as UTF-8 text."""
    try:
        text = path.read_text(encoding="utf-8", errors="ignore").strip()
        return [{"page": 1, "text": text}] if text else []
    except Exception as e:
        logger.warning("Could not read %s as text: %s", path.name, e)
        return []


# Format dispatch table
_EXTRACTORS = {
    ".pdf":  extract_pdf,
    ".pptx": extract_pptx,
    ".ppt":  extract_pptx,
    ".docx": extract_docx,
    ".xlsx": extract_xlsx,
    ".xlsm": extract_xlsx,
    ".xlsb": extract_xlsx,
    ".html": extract_html,
    ".htm": extract_html,
}

# These are plain text — no special library needed
_TEXT_SUFFIXES = {
    ".md", ".txt", ".csv", ".json", ".yaml", ".yml",
    ".xml", ".rst", ".tex",
    ".py", ".js", ".ts", ".java", ".c", ".cpp", ".h",
    ".sh", ".bat", ".ps1", ".r", ".sql",
}


def extract_pages(path: Path) -> list[dict]:
    """Route a file to the right extractor. Returns [] if nothing to extract."""
    suffix = path.suffix.lower()

    if suffix in _SKIP_SUFFIXES or path.name.startswith("."):
        return []

    extractor = _EXTRACTORS.get(suffix)
    if extractor:
        try:
            return extractor(path)
        except Exception as e:
            logger.warning("Extractor failed for %s (%s), trying text fallback", path.name, e)
            return extract_text(path)

    if suffix in _TEXT_SUFFIXES:
        return extract_text(path)

    # Unknown type — try text, silently skip if it looks binary
    try:
        text = path.read_text(encoding="utf-8", errors="strict").strip()
        if text:
            logger.debug("Read unknown type %s as plain text", path.name)
            return [{"page": 1, "text": text}]
    except (UnicodeDecodeError, Exception):
        pass  # Binary file — skip silently

    return []


# ── Chunking ────────────────────────────────────────────────────────────────

def chunk_text(text: str, chunk_size: int = 256, overlap: int = 32) -> list[str]:
    """
    Split text into overlapping word-count chunks.
    512 words ≈ 600-700 tokens — safe for bge-m3's 8192 token limit.
    """
    words = text.split()
    chunks, start = [], 0
    while start < len(words):
        chunks.append(" ".join(words[start : start + chunk_size]))
        start += chunk_size - overlap
    return chunks


def _chunk_id(source: str, page: int, chunk_idx: int) -> str:
    return hashlib.md5(f"{source}::p{page}::c{chunk_idx}".encode()).hexdigest()


# ── Ingest pipeline ─────────────────────────────────────────────────────────

async def ingest_file(path: Path, collection: chromadb.Collection) -> int:
    """
    Ingest a single file. Returns number of chunks added.
    Uses a semaphore to cap concurrent embedding calls.
    """
    pages = extract_pages(path)
    if not pages:
        return 0

    source_name = path.name
    suffix = path.suffix.lower()

    all_chunks, all_ids, all_meta = [], [], []
    for page_data in pages:
        for idx, chunk in enumerate(chunk_text(page_data["text"])):
            chunk_text_with_header = (
                f"Source: {source_name} | Page {page_data['page']} | Chunk {idx + 1}\n"
                f"{chunk}"
            )
            all_chunks.append(chunk_text_with_header)
            all_ids.append(_chunk_id(source_name, page_data["page"], idx))
            all_meta.append({
                "source": source_name,
                "page": page_data["page"],
                "chunk_index": idx,
                "type": suffix.lstrip("."),
            })

    if not all_chunks:
        return 0
    
    print(f"📄 {source_name}: {len(pages)} page(s) → {len(all_chunks)} chunk(s)")

    EMBED_BATCH = 16
    async with _EMBED_SEMAPHORE:
        for i in range(0, len(all_chunks), EMBED_BATCH):
            b_texts = all_chunks[i : i + EMBED_BATCH]
            b_ids   = all_ids[i : i + EMBED_BATCH]
            b_meta  = all_meta[i : i + EMBED_BATCH]

            try:
                vectors = await embed(b_texts)
            except Exception as e:
                logger.error("Embedding failed for %s batch %d: %s — skipping batch", source_name, i, e)
                continue  # skip this batch, keep going with the rest of the file

            collection.upsert(ids=b_ids, embeddings=vectors, documents=b_texts, metadatas=b_meta)

    logger.info("Ingested %-40s  %3d chunks from %d page(s)", source_name, len(all_chunks), len(pages))
    print(f"✅ {source_name}: {len(all_chunks)} chunks")
    return len(all_chunks)


async def ingest_all(material_path: str = config.COURSE_MATERIAL_PATH) -> dict:
    """
    Ingest all supported files from the course material directory.
    Returns {"files": N, "chunks": N, "skipped": N}.
    """
    path = Path(material_path)
    if not path.exists():
        path.mkdir(parents=True)
        logger.info("Created course material directory: %s", path)

    all_files = [p for p in path.rglob("*") if p.is_file()]
    if not all_files:
        logger.warning("No files found in %s", path)
        return {"files": 0, "chunks": 0, "skipped": 0}

    collection = await ensure_collection_compatible(get_collection())
    total_chunks = 0
    skipped = 0

    for f in tqdm(all_files, desc="Ingesting course material"):
        try:
            chunks = await ingest_file(f, collection)
            if chunks == 0:
                skipped += 1
            else:
                total_chunks += chunks
        except Exception as e:
            logger.error("Unexpected error ingesting %s: %s", f.name, e)
            skipped += 1

    ingested_files = len(all_files) - skipped
    summary = {"files": ingested_files, "chunks": total_chunks, "skipped": skipped}
    logger.info("Ingestion complete: %s", summary)
    return summary


def list_sources() -> list[str]:
    """Return a deduplicated list of ingested source document names."""
    collection = get_collection()
    results = collection.get(include=["metadatas"])
    if not results or not results.get("metadatas"):
        return []
    return sorted({m["source"] for m in results["metadatas"] if m})